import asyncio
import logging
import os
import uuid
from typing import List, Optional

import openai
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from hub.api.v1.models import FILE_URI_PREFIX, S3_URI_PREFIX
from hub.api.v1.sql import SqlClient, VectorStoreFile

logger = logging.getLogger(__name__)

# Constants for chunking strategy
CHARS_PER_TOKEN = 4
TOKEN_LIMIT = 800
CHUNK_SIZE = TOKEN_LIMIT * CHARS_PER_TOKEN
CHUNK_OVERLAP = CHUNK_SIZE // 4

# File extensions
OFFICE_EXTENSIONS = [".docx", ".pptx", ".xlsx"]

# Embedding model
EMBEDDING_MODEL = "nomic-ai/nomic-embed-text-v1.5"

"""
Chunking strategy:
- CHARS_PER_TOKEN: Approximate average number of characters per token.
- TOKEN_LIMIT: Approximate token limit for the embedding model.
- CHUNK_SIZE: Set to TOKEN_LIMIT * CHARS_PER_TOKEN to create chunks that fit within the model's token limit.
- CHUNK_OVERLAP: Set to 1/4 of CHUNK_SIZE to ensure context continuity between chunks.

This strategy creates chunks based on the estimated token limit of the embedding model.
The overlap ensures smooth transitions and context preservation between chunks.
These values may need adjustment based on the specific requirements of the embedding model in use.
"""


async def generate_embeddings_for_vector_store(vector_store_id: str):
    """Generate embeddings for all files in a vector store.

    Args:
    ----
        vector_store_id (str): The ID of the vector store.

    Raises:
    ------
        ValueError: If the vector store with the given ID is not found.

    """
    logger.info(f"Starting embedding generation for vector store: {vector_store_id}")
    sql_client = SqlClient()
    vector_store = sql_client.get_vector_store(vector_store_id=vector_store_id)

    if not vector_store:
        logger.error(f"Vector store with id {vector_store_id} not found")
        raise ValueError(f"Vector store with id {vector_store_id} not found")

    tasks = [
        generate_embeddings_for_file(file_id, vector_store_id, vector_store.chunking_strategy)
        for file_id in vector_store.file_ids
    ]
    await asyncio.gather(*tasks)

    logger.info(f"Finished embedding generation for vector store: {vector_store_id}")


async def generate_embeddings_for_file(file_id: str, vector_store_id: str, chunking_strategy: Optional[dict] = None):
    """Generate embeddings for a specific file and store them in the vector store.

    Args:
    ----
        file_id (str): The ID of the file to generate embeddings for.
        vector_store_id (str): The ID of the vector store to associate the embeddings with.
        chunking_strategy (dict, optional): Chunking strategy to use for splitting the file content.

    Raises:
    ------
        ValueError: If the file with the given ID is not found.

    """
    logger.info(f"Starting embedding generation for file: {file_id}")

    sql_client = SqlClient()
    file_details = sql_client.get_file_details(file_id)
    if not file_details:
        logger.error(f"File with id {file_id} not found")
        raise ValueError(f"File with id {file_id} not found")

    content = await get_file_content(file_details)
    chunks = create_chunks(content, chunking_strategy)
    logger.debug(f"Created {len(chunks)} chunks for file: {file_id}")

    embedding_tasks = [generate_embedding(chunk) for chunk in chunks]
    embeddings = await asyncio.gather(*embedding_tasks)

    for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
        embedding_id = f"vfe_{uuid.uuid4().hex[:24]}"
        try:
            sql_client.store_embedding(
                id=embedding_id,
                vector_store_id=vector_store_id,
                file_id=file_id,
                chunk_index=i,
                chunk_text=chunk,
                embedding=embedding,
            )
        except Exception as e:
            logger.error(f"Failed to store embedding: {embedding_id} for file: {file_id}, error: {e}")

    sql_client.update_file_embedding_status(file_id, "completed")

    if embeddings:
        embedding_dimensions = len(embeddings[0])
        sql_client.update_vector_store_embedding_info(vector_store_id, EMBEDDING_MODEL, embedding_dimensions)

    logger.info(f"Finished embedding generation for file: {file_id}")


def create_chunks(text: str, chunking_strategy=None) -> List[str]:
    """Split the input text into chunks of appropriate size for embedding generation.

    Args:
    ----
        text (str): The input text to be split into chunks.
        chunking_strategy (ChunkingStrategy): Optional strategy to use for splitting the file content.

    Returns:
    -------
        List[str]: A list of text chunks.

    """
    chunk_size = CHUNK_SIZE
    chunk_overlap = CHUNK_OVERLAP
    if chunking_strategy:
        chunk_size = chunking_strategy.get("max_chunk_size_tokens", CHUNK_SIZE)
        chunk_overlap = chunking_strategy.get("chunk_overlap_tokens", CHUNK_OVERLAP)

    chunks = recursive_split(text, chunk_size, chunk_overlap)
    logger.debug(f"Created {len(chunks)} chunks, sizes: {[len(chunk) for chunk in chunks]}")
    return chunks


def recursive_split(text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
    """Recursively split text into chunks of specified size with overlap.

    This function attempts to split the text using various separators, starting with
    paragraph breaks and moving to smaller separators if needed. It ensures that
    chunks do not exceed the specified size and maintains the required overlap.

    Args:
    ----
        text (str): Input text to split.
        chunk_size (int): Maximum size of each chunk.
        chunk_overlap (int): Overlap size between chunks.

    Returns:
    -------
        List[str]: List of text chunks.

    """
    logger.debug(f"Splitting text into chunks of size {chunk_size} with overlap {chunk_overlap}, length: {len(text)}")
    if len(text) <= chunk_size:
        return [text]

    separators = ["\n\n", "\n", ". ", " ", ""]
    for separator in separators:
        if separator:
            splits = text.split(separator)
        else:
            splits = list(text)

        chunks: List[str] = []
        current_chunk = ""

        for split in splits:
            split_with_sep = split + (separator if separator else "")

            if len(current_chunk) + len(split_with_sep) > chunk_size:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = split_with_sep
            else:
                current_chunk += split_with_sep

            while len(current_chunk) > chunk_size:
                chunks.append(current_chunk[:chunk_size].strip())
                current_chunk = current_chunk[chunk_size - chunk_overlap :]

        if current_chunk:
            chunks.append(current_chunk.strip())

        if len(chunks) > 1:
            return chunks

    # If we reach here, it means we couldn't split the text
    return [text[:chunk_size]]


async def generate_embedding(text: str, query: bool = False):
    """Generate an embedding for the given text using the Nomic AI model.

    Args:
    ----
        text (str): The text to generate an embedding for.
        query (bool, optional): If True, the text is treated as a search query.
            If False, it's treated as a document. Defaults to False.

    Returns:
    -------
        list: The embedding vector for the input text.

    """
    client = openai.AsyncOpenAI(
        base_url="https://api.fireworks.ai/inference/v1", api_key=os.getenv("FIREWORKS_API_KEY")
    )
    prefix = "search_query: " if query else "search_document: "
    response = await client.embeddings.create(input=prefix + text, model=EMBEDDING_MODEL)
    return response.data[0].embedding


async def get_file_content(file_details: VectorStoreFile) -> str:
    """Retrieve the content of a file based on its URI.

    This function supports both local file system and S3 storage. For S3 files,
    it downloads the file to a temporary location before extracting the content.

    Args:
    ----
        file_details (VectorStoreFile): Details of the file, including its URI and encoding.

    Returns:
    -------
        str: The content of the file.

    Raises:
    ------
        ValueError: If the file URI is not supported.

    """
    logger.info(f"Getting content for file: {file_details.file_uri}")
    encoding = file_details.encoding or "utf-8"

    if file_details.file_uri.startswith(FILE_URI_PREFIX):
        file_path = file_details.file_uri[len(FILE_URI_PREFIX) :]
        logger.debug(f"Extracting content from local file: {file_path}")
        return extract_content(file_path, encoding)
    elif file_details.file_uri.startswith(S3_URI_PREFIX):
        logger.debug(f"Extracting content from S3 file: {file_details.file_uri}")
        import boto3

        s3_client = boto3.client("s3")
        parts = file_details.file_uri[len(S3_URI_PREFIX) :].split("/", 1)
        if len(parts) != 2:
            logger.error(f"Invalid S3 URI format: {file_details.file_uri}")
            raise ValueError(f"Invalid S3 URI format: {file_details.file_uri}")
        bucket_name, key = parts
        response = s3_client.get_object(Bucket=bucket_name, Key=key)
        temp_file_path = f"/tmp/tempfile_{uuid.uuid4().hex}_{file_details.filename}"
        with open(temp_file_path, "wb") as f:
            f.write(response["Body"].read())
        logger.debug(f"Downloaded S3 file to temporary path: {temp_file_path}")
        content = extract_content(temp_file_path, encoding)
        os.remove(temp_file_path)
        logger.debug(f"Removed temporary file: {temp_file_path}")
        return content
    else:
        logger.error(f"Unsupported file URI: {file_details.file_uri}")
        raise ValueError(f"Unsupported file URI: {file_details.file_uri}")


def extract_content(file_path: str, encoding: str = "utf-8") -> str:
    """Extract content from various file types.

    This function supports PDF, DOCX, PPTX, XLSX, and plain text files.
    It determines the file type based on the file extension and uses
    the appropriate extraction method.

    Args:
    ----
        file_path (str): Path to the file.
        encoding (str, optional): Encoding for text files. Defaults to "utf-8".

    Returns:
    -------
        str: The extracted content of the file.

    """
    logger.debug(f"Extracting content from file: {file_path}")
    _, file_extension = os.path.splitext(file_path.lower())

    if file_extension == ".pdf":
        logger.debug("Detected PDF file, using PDF extraction method")
        return extract_pdf_content(file_path)
    elif file_extension == ".docx":
        logger.debug("Detected DOCX file, using python-docx extraction method")
        return extract_docx_content(file_path)
    elif file_extension == ".pptx":
        logger.debug("Detected PPTX file, using python-pptx extraction method")
        return extract_pptx_content(file_path)
    elif file_extension == ".xlsx":
        logger.debug("Detected XLSX file, using openpyxl extraction method")
        return extract_xlsx_content(file_path)
    else:
        logger.debug("Detected text file, using standard text extraction method")
        return extract_text_file(file_path, encoding)


def extract_text_file(file_path: str, encoding: str) -> str:
    logger.debug(f"Extracting content from text file: {file_path}")
    try:
        with open(file_path, "r", encoding=encoding) as file:
            content = file.read()
        logger.debug(f"Successfully extracted content from text file: {file_path}")
        return content
    except UnicodeDecodeError:
        logger.error(f"Unable to decode {file_path} with encoding {encoding}")
        return f"Error: Unable to decode {file_path}"


def extract_pdf_content(file_path: str) -> str:
    logger.debug(f"Extracting content from PDF file: {file_path}")
    with open(file_path, "rb") as file:
        reader = PdfReader(file)
        content = "\n".join(page.extract_text() for page in reader.pages)
    logger.debug(f"Successfully extracted content from PDF file: {file_path}")
    return content


def extract_docx_content(file_path: str) -> str:
    logger.debug(f"Extracting content from DOCX file: {file_path}")
    doc = Document(file_path)
    content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    logger.debug(f"Successfully extracted content from DOCX file: {file_path}")
    return content


def extract_pptx_content(file_path: str) -> str:
    logger.debug(f"Extracting content from PPTX file: {file_path}")
    prs = Presentation(file_path)
    content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)
    logger.debug(f"Successfully extracted content from PPTX file: {file_path}")
    return "\n".join(content)


def extract_xlsx_content(file_path: str) -> str:
    logger.debug(f"Extracting content from XLSX file: {file_path}")
    wb = load_workbook(file_path, read_only=True)
    content = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            content.append("\t".join(str(cell) for cell in row if cell is not None))
    logger.debug(f"Successfully extracted content from XLSX file: {file_path}")
    return "\n".join(content)
